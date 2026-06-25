"""Build the pre-defense pptx deck (redesigned 2026-05-13).

Output:
- presentation/Фамилия_Имя_Предзащита.pptx   — 14 numbered + 5 annex slides
  Re-running this script overwrites the file in place; manual formatting in
  PowerPoint is lost. Run from repo root: ``python scripts/build_predefense_pptx.py``.

The deck is aligned with presentation/ch1_full.md (Chapter 1 of the diploma)
and standard Russian diploma-defense criteria. Two slides are drafted from
the user's message of 2026-05-13 and clearly marked
[USER DRAFT — please revise]:
  - Slide 5:  Architecture
  - Slide 6:  Metrics

Placeholders [ФИО] / [РУКОВОДИТЕЛЬ] / [ВУЗ / ПРОГРАММА] / [ДАТА] are intentionally
literal so the author can substitute them in PowerPoint with Find & Replace, or
rerun this script after editing the SUBSTITUTIONS dict below.
"""

from __future__ import annotations

from pathlib import Path
from typing import Iterable

import fitz
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Pt

# Configuration -------------------------------------------------------------
ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "Yandex_template_Arial_ppt.pptx"
OUT_DIR = ROOT / "presentation"
FIG_DIR = ROOT / "experiments" / "baseline" / "figures_real_augmented_best"

SUBSTITUTIONS = {
    "FIO": "[ФИО]",
    "SUPERVISOR": "[РУКОВОДИТЕЛЬ]",
    "UNIVERSITY": "[ВУЗ / факультет / программа]",
    "DATE": "[ДАТА]",
    "EMAIL": "[email]",
}

# 13 numbered content slides; no backup slides in this revision.
TOTAL = 13

# Cache for PDF→PNG conversions
PNG_CACHE = ROOT / ".cache" / "predefense_pngs"
PNG_CACHE.mkdir(parents=True, exist_ok=True)


# Helpers -------------------------------------------------------------------
def pdf_first_page_to_png(pdf_path: Path, zoom: float = 2.5) -> Path:
    """Rasterise page 1 of a PDF to PNG (cached). Returns the PNG path."""
    pdf_path = Path(pdf_path)
    stem = pdf_path.with_suffix("").as_posix().replace("/", "__")
    out = PNG_CACHE / f"{Path(stem).name}_{zoom:.1f}.png"
    if out.exists():
        return out
    doc = fitz.open(pdf_path)
    pix = doc[0].get_pixmap(matrix=fitz.Matrix(zoom, zoom))
    pix.save(out.as_posix())
    doc.close()
    return out


def emu(inches: float) -> int:
    return int(inches * 914400)


def add_footer(slide, n: int, total: int, slide_w: int, slide_h: int) -> None:
    """Bottom-right N/Total footer."""
    tb = slide.shapes.add_textbox(
        left=slide_w - emu(2.0),
        top=slide_h - emu(0.7),
        width=emu(1.8),
        height=emu(0.5),
    )
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{n}/{total}"
    run.font.name = "Arial"
    run.font.size = Pt(20)
    run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)


def add_text(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text_blocks: Iterable[tuple[str, int, bool]],
    *,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    color: RGBColor | None = None,
) -> None:
    """Each block is (text, point_size, bold). Empty text => blank line."""
    tb = slide.shapes.add_textbox(emu(left), emu(top), emu(width), emu(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = emu(0.1)
    tf.margin_right = emu(0.1)
    tf.margin_top = emu(0.05)
    tf.margin_bottom = emu(0.05)
    first = True
    for text, size, bold in text_blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.name = "Arial"
        run.font.size = Pt(size)
        run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color


def add_bullets(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    bullets: Iterable[tuple[int, str, int]],
    *,
    color: RGBColor | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    """Each bullet is (level, text, point_size). level=0 is a top-level bullet."""
    tb = slide.shapes.add_textbox(emu(left), emu(top), emu(width), emu(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = emu(0.1)
    tf.margin_right = emu(0.1)
    tf.margin_top = emu(0.05)
    tf.margin_bottom = emu(0.05)
    first = True
    for level, text, size in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.alignment = align
        prefix = "• " if level == 0 else ("— " if level == 1 else "· ")
        run = p.add_run()
        run.text = prefix + text
        run.font.name = "Arial"
        run.font.size = Pt(size)
        if color is not None:
            run.font.color.rgb = color


def add_image(slide, png_path: Path, left: float, top: float, width: float, height: float | None = None) -> None:
    """Insert an image keeping aspect ratio inside the (width, height) box."""
    if height is None:
        slide.shapes.add_picture(str(png_path), emu(left), emu(top), width=emu(width))
        return
    with Image.open(png_path) as im:
        iw, ih = im.size
    img_ar = iw / ih
    box_ar = width / height
    if img_ar > box_ar:
        w_emu = emu(width)
        h_emu = int(emu(width) / img_ar)
        x = emu(left)
        y = emu(top) + (emu(height) - h_emu) // 2
    else:
        h_emu = emu(height)
        w_emu = int(emu(height) * img_ar)
        x = emu(left) + (emu(width) - w_emu) // 2
        y = emu(top)
    slide.shapes.add_picture(str(png_path), x, y, width=w_emu, height=h_emu)


def add_notes(slide, text: str) -> None:
    ns = slide.notes_slide
    tf = ns.notes_text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.strip()
    run.font.name = "Arial"
    run.font.size = Pt(12)


def add_rect(slide, left: float, top: float, width: float, height: float, fill: RGBColor):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, emu(left), emu(top), emu(width), emu(height)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.fill.background()
    return shp


def add_title(slide, text: str, *, top: float = 0.5, height: float = 1.4, size: int = 44, left: float = 0.8, width: float = 24.5) -> None:
    add_text(
        slide,
        left=left, top=top, width=width, height=height,
        text_blocks=[(text, size, True)],
        align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,
    )


def add_draft_tag(slide, top: float = 14.2) -> None:
    """Visible '[draft, will be revised]' tag at the bottom-left of a slide."""
    add_text(
        slide,
        left=0.8, top=top, width=18.0, height=0.4,
        text_blocks=[("[черновик автора — будет переписан]", 14, False)],
        color=RGBColor(0x99, 0x99, 0x99),
    )


def add_arrow(slide, left: float, top: float, width: float = 0.6, height: float = 0.3, fill: RGBColor | None = None):
    """Right-pointing arrow shape."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, emu(left), emu(top), emu(width), emu(height)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill if fill is not None else RGBColor(0x55, 0x55, 0x55)
    shp.line.fill.background()
    return shp


def add_table(slide, left: float, top: float, headers, rows, col_widths, *, header_size: int = 20, body_size: int = 18, row_h: float = 0.7, highlight=None):
    """Convenience table builder."""
    n_cols = len(headers)
    n_rows = 1 + len(rows)
    col_widths_emu = [emu(w) for w in col_widths]
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, emu(left), emu(top),
                                       sum(col_widths_emu), emu(row_h) * n_rows)
    table = tbl_shape.table
    for i, w in enumerate(col_widths_emu):
        table.columns[i].width = w
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = "Arial"
                r.font.size = Pt(header_size)
                r.font.bold = True
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = "Arial"
                    r.font.size = Pt(body_size)
                    if highlight is not None and val == highlight:
                        r.font.bold = True
                        r.font.color.rgb = RGBColor(0xC0, 0x10, 0x10)
    return tbl_shape


# Build ---------------------------------------------------------------------
def build() -> Path:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(TEMPLATE))

    # Drop the template's starter slide(s) cleanly so we do not leave orphan parts in the zip.
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
        if rId:
            prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    blank = prs.slide_layouts[43]  # 'Пустой слайд' — fully empty
    sw = prs.slide_width
    sh = prs.slide_height

    # 7 numbered slides. Some legacy function names are kept; the visible
    # slide number is controlled by the add_footer(...) call inside each.
    _slide_1_title(prs, blank, sw, sh)
    _slide_2_actuality(prs, blank, sw, sh)
    _slide_3_goals_novelty(prs, blank, sw, sh)
    _slide_4_vsh_setup(prs, blank, sw, sh)            # merged VSH + E→P collapse
    _slide_8_architecture_DRAFT(prs, blank, sw, sh)   # slide 6/7
    _slide_9_metrics_DRAFT(prs, blank, sw, sh)        # slide 6/13
    _slide_run_best_small(prs, blank, sw, sh)         # slide 7/13
    _slide_run_coef(prs, blank, sw, sh)               # slide 8/13
    _slide_run_coef_raw_flat(prs, blank, sw, sh)      # slide 9/13
    _slide_run_physics(prs, blank, sw, sh)            # slide 10/13
    _slide_run_best(prs, blank, sw, sh)               # slide 11/13
    _slide_holdout_field_grid(prs, blank, sw, sh)     # slide 12/13
    _slide_conclusions(prs, blank, sw, sh)            # slide 13/13

    out_path = OUT_DIR / "Фамилия_Имя_Предзащита.pptx"
    prs.save(str(out_path))
    return out_path


# ----------------------------------------------------------------------
# SLIDE 1 — TITLE
# ----------------------------------------------------------------------
def _slide_1_title(prs, blank, sw, sh) -> None:
    s = prs.slides.add_slide(blank)
    add_rect(s, 0.0, 0.0, 26.66, 0.4, RGBColor(0xFF, 0x00, 0x00))
    add_text(
        s, 0.8, 1.8, 24.5, 1.4,
        text_blocks=[("Предзащита диплома", 28, False)],
        anchor=MSO_ANCHOR.TOP,
    )
    add_text(
        s, 0.8, 3.0, 24.5, 5.0,
        text_blocks=[
            ("Мультипольное разложение", 60, True),
            ("электромагнитного поля", 60, True),
            ("методами машинного обучения", 60, True),
        ],
        anchor=MSO_ANCHOR.TOP,
    )
    add_text(
        s, 0.8, 9.8, 24.5, 4.5,
        text_blocks=[
            (f"Автор: {SUBSTITUTIONS['FIO']}", 28, False),
            (f"Научный руководитель: {SUBSTITUTIONS['SUPERVISOR']}", 28, False),
            (SUBSTITUTIONS["UNIVERSITY"], 28, False),
            ("", 12, False),
            (SUBSTITUTIONS["DATE"], 24, False),
        ],
    )
    add_footer(s, 1, TOTAL, sw, sh)
    add_notes(
        s,
        "Добрый день. Меня зовут [ФИО]. Сегодня я представляю предзащиту дипломной работы "
        "«Мультипольное разложение электромагнитного поля методами машинного обучения». "
        "Научный руководитель — [РУКОВОДИТЕЛЬ]. Доклад займёт около десяти минут, после "
        "чего я готова ответить на вопросы. Главная идея: мы строим обученную модель, "
        "которая по измеренной мощности дальнего поля антенны восстанавливает мультипольные "
        "коэффициенты — компактную физическую подпись радиатора. Аналитически в безфазовом "
        "режиме измерения это сделать нельзя — мы покажем это на слайде 4."
    )


# ----------------------------------------------------------------------
# SLIDE 2 — Актуальность
# ----------------------------------------------------------------------
def _slide_2_actuality(prs, blank, sw, sh) -> None:
    s = prs.slides.add_slide(blank)
    add_title(s, "Актуальность: безфазовое измерение и физическая подпись радиатора")
    add_bullets(
        s, 0.8, 2.2, 14.5, 11.0,
        bullets=[
            (0, "Промышленные измерения дальнего поля антенн — безфазовые: на сетке 360 × 179 направлений измеряется только мощность P (Hansen 1988).", 24),
            (0, "Когерентные комплексные измерения (амплитуда + фаза) требуют фазостабильной аппаратуры, дороги и медленны — стандартны только в метрологии.", 24),
            (0, "Мультипольные коэффициенты (aᴱ, aᴹ) — каноническая компактная подпись радиатора (Jackson 1999, Гл. 9).", 24),
            (0, "При L = 15: 1 020 вещественных коэффициентов против 64 440 пикселей P → сжатие ≈ 250 : 1.", 24),
            (0, "Аналитический инверс существует только для комплексного поля E. Из P он структурно невозможен (доказательство — слайды 5–6).", 24),
        ],
    )

    # Right-side ratio graphic — three coloured bars showing relative DoF.
    add_rect(s, 16.0, 2.5, 10.5, 10.5, RGBColor(0xF7, 0xF7, 0xF7))
    add_text(
        s, 16.3, 2.8, 9.9, 0.8,
        text_blocks=[("Размерный баланс одного образца", 22, True)],
    )
    # Bar 1: full E (full width).
    add_rect(s, 16.3, 4.0, 9.9, 1.4, RGBColor(0xCC, 0xDD, 0xFF))
    add_text(
        s, 16.5, 4.0, 9.5, 1.4,
        text_blocks=[
            ("Комплексное поле E (4 числа / точка)", 18, True),
            ("257 760 ℝ-чисел", 22, False),
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # Bar 2: P — 1/4 of E (proportional width, 9.9 × 0.25 = 2.475).
    add_rect(s, 16.3, 5.7, 2.5, 1.4, RGBColor(0xFF, 0xCC, 0xCC))
    add_text(
        s, 16.5, 5.7, 9.5, 1.4,
        text_blocks=[
            ("Мощность P (1 число / точка)", 18, True),
            ("64 440 ℝ-чисел  (¾ информации потеряно)", 22, False),
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # Bar 3: coefficients — width ≈ 9.9 × 1020 / 257760 ≈ 0.039 (effectively a sliver).
    add_rect(s, 16.3, 7.4, 0.4, 1.4, RGBColor(0xFF, 0xEE, 0x99))
    add_text(
        s, 16.5, 7.4, 9.5, 1.4,
        text_blocks=[
            ("Коэффициенты (aᴱ, aᴹ)", 18, True),
            ("1 020 ℝ-чисел  (×250 сжатие)", 22, False),
        ],
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        s, 16.3, 9.5, 9.9, 3.0,
        text_blocks=[
            ("Цель работы:", 22, True),
            ("восстановить нижнее представление", 20, False),
            ("из среднего, не зная верхнее.", 20, False),
        ],
    )
    add_footer(s, 2, TOTAL, sw, sh)
    add_notes(
        s,
        "Мотивация. Антенну полностью описывает комплексное дальнее поле на сфере — это 257 760 "
        "вещественных чисел. Промышленные измерения безфазовые: измеряется только мощность, это "
        "уже 64 440. Три четверти угловой информации отсутствуют физически. Мультипольные "
        "коэффициенты — всего 1 020 чисел: компактная и физически осмысленная подпись радиатора "
        "в его собственном базисе. Аналитический инверс существует только для комплексного поля; "
        "из мощности он структурно невозможен. Эта неустранимость и мотивирует обращаться к "
        "обучаемой модели."
    )


# ----------------------------------------------------------------------
# SLIDE 3 — Цель, задачи, новизна
# ----------------------------------------------------------------------
def _slide_3_goals_novelty(prs, blank, sw, sh) -> None:
    s = prs.slides.add_slide(blank)
    add_title(s, "Цель работы, задачи, научная новизна")

    # Left column — goal and tasks.
    add_text(
        s, 0.8, 2.2, 13.0, 1.0,
        text_blocks=[("Цель и задачи", 30, True)],
    )
    add_text(
        s, 0.8, 3.2, 13.0, 1.6,
        text_blocks=[(
            "Построить обучаемую модель-регуляризатор f_η: P → (aᴱ, aᴹ), решающую безфазовый VSH-инверс на реальных антеннах.",
            20, False
        )],
    )
    add_bullets(
        s, 0.8, 5.0, 13.0, 8.5,
        bullets=[
            (0, "Г1. Формализовать обратную задачу P → (aᴱ, aᴹ) и рассказать о неоднозначности её решения.", 20),
            (0, "Г2. Создать фреймворк для гибкого обучения моделей в задачах мультипольного анализа.", 20),
            (0, "Г3. Превзойти тривиальный нейросетевой бейзлайн на реальных данных.", 20),
            (0, "Г4. Исследовать архитектурные возможности и ограничения в задаче.", 20),
            (0, "Г5. Численно подтвердить полезность обучаемого приора на реальной отложенной выборке.", 20),
        ],
    )

    # Right column — novelty.
    add_text(
        s, 14.2, 2.2, 12.3, 1.0,
        text_blocks=[("Научная новизна", 30, True)],
    )
    add_bullets(
        s, 14.2, 3.2, 12.3, 10.5,
        bullets=[
            (0, "Гибкий фреймворк сквозного обучения для задачи мультипольного анализа: расширяемые feature-pipelines, лосс-блоки и аугментации, физически согласованные с VSH-разложением.", 20),
            (0, "Применение fully-learned baseline (Arridge et al. 2019, §5.1.3) к VSH-power инверсу на реальных антеннах с регуляризацией power_loss с p_rank.", 20),
            (0, "Эмпирическое наблюдение: главный рычаг качества — распределение обучающих данных, а не архитектура (≈ 100× эффект режима генератора).", 20),
        ],
    )
    add_footer(s, 3, TOTAL, sw, sh)
    add_notes(
        s,
        "Цель и пять задач. Г1 — формализация и описание неоднозначности решения; Г2 — гибкий "
        "фреймворк обучения для задач мультипольного анализа; Г3 — превзойти тривиальный "
        "нейросетевой бейзлайн на реальных данных; Г4 — исследование архитектурных возможностей "
        "и ограничений; Г5 — численное подтверждение полезности обучаемого приора на реальной "
        "отложенной выборке. Научная новизна — три пункта: гибкий фреймворк сквозного обучения, "
        "применение fully-learned подхода (Arridge et al. 2019) к VSH-инверсу на реальных "
        "антеннах с регуляризацией p_rank, и эмпирический вывод о доминировании распределения "
        "обучающих данных над архитектурой."
    )


# ----------------------------------------------------------------------
# SLIDE 4 — Физическая постановка: VSH
# ----------------------------------------------------------------------
def _slide_4_vsh_setup(prs, blank, sw, sh) -> None:
    """Merged: VSH formal setup + analytic inverse + E→P collapse."""
    s = prs.slides.add_slide(blank)
    add_title(s, "Физика задачи: VSH-разложение и коллапс E → P")

    # VSH expansion formula at the top, full-width.
    add_rect(s, 0.8, 1.9, 25.5, 1.8, RGBColor(0xF7, 0xF7, 0xF7))
    add_text(
        s, 0.8, 2.0, 25.5, 0.7,
        text_blocks=[("VSH-разложение поля дальней зоны (Jackson 1999, Гл. 9):", 20, True)],
    )
    add_text(
        s, 0.8, 2.7, 25.5, 1.0,
        text_blocks=[
            ("E(θ, φ) = Σ_l Σ_m [ aᴱₗₘ · Ψᴱₗₘ(θ, φ)  +  aᴹₗₘ · Ψᴹₗₘ(θ, φ) ]", 24, False),
        ],
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )

    # Left column — VSH structure + grid.
    add_text(
        s, 0.8, 4.0, 12.5, 1.0,
        text_blocks=[("Структура базиса и сетки наблюдения", 24, True)],
    )
    add_bullets(
        s, 0.8, 5.0, 12.5, 8.0,
        bullets=[
            (0, "Две семьи мод: электрические (TM) и магнитные (TE).", 22),
            (0, "Базис ортонормирован под area-weighted L²(S²): ⟨Ψˣₗₘ, Ψˣ'ₗ'ₘ'⟩ = δ_XX' δ_ll' δ_mm'.", 22),
            (0, "Усечение L = 15 → K = L(L+2) = 255 мод/семья.", 22),
            (0, "Упаковка [Re aᴱ, Im aᴱ, Re aᴹ, Im aᴹ] длиной 4K = 1 020.", 22),
            (0, "Сетка 360 × 179: θ ∈ [1°, 179°], φ ∈ [0°, 359°], шаг 1°. Полюса исключены (sin θ = 0).", 22),
            (0, "Один образец поля: 4 × 360 × 179 = 257 760 ℝ-DoF.", 22),
        ],
    )

    # Right column — analytic inversion (top half) and the collapse (bottom half).
    add_rect(s, 13.8, 4.0, 12.7, 4.4, RGBColor(0xEE, 0xF7, 0xEE))
    add_text(
        s, 14.0, 4.1, 12.3, 0.9,
        text_blocks=[("Когда известен E:  gold standard", 22, True)],
        color=RGBColor(0x22, 0x66, 0x22),
    )
    add_bullets(
        s, 14.0, 5.0, 12.3, 3.4,
        bullets=[
            (0, "Ортогональность ⇒ aˣₗₘ = ⟨E_UT, Ψˣₗₘ⟩_{S²}: один inner-product на коэффициент.", 20),
            (0, "В дискрете: â = B · vec(E_UT). Детерминирован, линеен, точен.", 20),
            (0, "Эталон, против которого сравниваем обучаемые модели.", 20),
        ],
    )

    add_rect(s, 13.8, 8.6, 12.7, 4.7, RGBColor(0xFF, 0xF0, 0xEE))
    add_text(
        s, 14.0, 8.7, 12.3, 0.9,
        text_blocks=[("Когда доступен только P:  4 → 1 коллапс", 22, True)],
        color=RGBColor(0xC0, 0x10, 0x10),
    )
    add_bullets(
        s, 14.0, 9.6, 12.3, 3.6,
        bullets=[
            (0, "P = |Eθ|² + |Eφ|²: 4 числа на точку → 1 число.", 20),
            (0, "Прообраз — 3-сфера S³ радиуса √P в ℝ⁴ (трёхпараметрическое семейство).", 20),
            (0, "Глобально: 257 760 → 64 440 ℝ-DoF, ¾ информации физически отсутствуют.", 20),
            (0, "Композитный форвард A = |·|² ∘ S не имеет аналитического обратного.", 20),
            (0, "Решений много — алгоритм может только выбрать одно.", 20),
        ],
    )

    add_footer(s, 4, TOTAL, sw, sh)
    add_notes(
        s,
        "Физика и формальная постановка одним слайдом. Сверху — VSH-разложение (Jackson 1999, "
        "Гл. 9): любое поле дальней зоны однозначно раскладывается по двум семьям векторных "
        "сферических гармоник — TM и TE. Слева — параметры базиса и сетки: L = 15, 255 мод "
        "на семью, упакованный вектор длиной 1 020, сетка 360 × 179, полюса исключены. "
        "Справа сверху — gold standard: при известном комплексном поле коэффициенты "
        "восстанавливаются одним матричным умножением. Справа снизу — что происходит, когда "
        "комплексное поле недоступно: переход от поля к мощности — поточечный коллапс четыре "
        "в один; глобально мы теряем три четверти угловой информации. Никакое матричное "
        "умножение не вернёт уничтоженную фазу — отсюда необходимость регуляризации."
    )



# ======================================================================
# SLIDE 5 — Архитектура
# ======================================================================
def _slide_8_architecture_DRAFT(prs, blank, sw, sh) -> None:
    s = prs.slides.add_slide(blank)
    add_title(s, "Архитектура")

    # Top — three columns of bullets: model, regularisation losses, features.
    add_text(
        s, 0.8, 2.0, 8.5, 1.0,
        text_blocks=[("Модель", 28, True)],
    )
    add_bullets(
        s, 0.8, 3.0, 8.5, 4.5,
        bullets=[
            (0, "Бейзлайн: MLP 5×200.", 22),
            (1, "5 скрытых слоёв × 200 нейронов.", 20),
            (1, "Без свёрток и attention — простой baseline.", 20),
            (1, "Кастомный Trainer без сторонних обучающих фреймворков.", 20),
        ],
    )

    add_text(
        s, 9.5, 2.0, 8.5, 1.0,
        text_blocks=[("Функции потерь / регуляризация", 28, True)],
    )
    add_bullets(
        s, 9.5, 3.0, 8.5, 4.5,
        bullets=[
            (0, "coef_mse — прямая MSE на упакованный ŷ.", 22),
            (0, "power_loss — area-weighted MSE на P̂ через дифференцируемый VSH-декодер.", 22),
            (0, "physics_power_rank, rank_bin_p — rank-based варианты (слайд 7).", 22),
            (0, "physics_power_mixed — coef_mse + λ · physics_power.", 22),
        ],
    )

    add_text(
        s, 18.2, 2.0, 8.3, 1.0,
        text_blocks=[("Признаки (input features)", 28, True)],
    )
    add_bullets(
        s, 18.2, 3.0, 8.3, 4.5,
        bullets=[
            (0, "raw_flat — P целиком (64 440 чисел).", 22),
            (0, "power_pca — PCA(P) до 128 компонент.", 22),
            (0, "cv_only — FFT-radial + sh_power (≈21 чисел).", 22),
            (0, "raw_plus_sh — raw + sh_power, главный composite.", 22),
            (0, "subsample_stride4 — каждая 4-я угловая точка.", 22),
        ],
    )

    # Highlight the best feature set.
    add_text(
        s, 0.8, 7.6, 25.5, 0.7,
        text_blocks=[("На реальных данных лучший feature set — raw_plus_sh (raw + spherical-harmonic power).", 22, True)],
        color=RGBColor(0xC0, 0x10, 0x10),
    )

    # Block diagram of the end-to-end pipeline (lower half).
    add_text(
        s, 0.8, 8.5, 25.5, 0.7,
        text_blocks=[("Сквозной пайплайн обучения:", 22, True)],
    )

    # Box 1 — P features.
    add_rect(s, 1.2, 9.5, 3.5, 1.6, RGBColor(0xCC, 0xDD, 0xFF))
    add_text(
        s, 1.2, 9.5, 3.5, 1.6,
        text_blocks=[
            ("Входные признаки", 14, True),
            ("raw_plus_sh ∈ ℝ^{≈64500}", 16, False),
        ],
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )
    add_arrow(s, 4.9, 10.15, width=0.7, height=0.4)

    # Box 2 — MLP.
    add_rect(s, 5.8, 9.5, 4.0, 1.6, RGBColor(0xDD, 0xEE, 0xCC))
    add_text(
        s, 5.8, 9.5, 4.0, 1.6,
        text_blocks=[
            ("MLP 5×200", 18, True),
            ("(параметры η)", 14, False),
        ],
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )
    add_arrow(s, 10.0, 10.15, width=0.7, height=0.4)

    # Box 3 — packed coefficients ŷ.
    add_rect(s, 10.9, 9.5, 4.0, 1.6, RGBColor(0xFF, 0xEE, 0x99))
    add_text(
        s, 10.9, 9.5, 4.0, 1.6,
        text_blocks=[
            ("ŷ — коэффициенты", 14, True),
            ("ℝ¹⁰²⁰ (packed)", 16, False),
        ],
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )
    add_arrow(s, 15.1, 10.15, width=0.7, height=0.4)

    # Box 4 — P-head (diff VSH decoder).
    add_rect(s, 16.0, 9.5, 4.0, 1.6, RGBColor(0xFF, 0xCC, 0xCC))
    add_text(
        s, 16.0, 9.5, 4.0, 1.6,
        text_blocks=[
            ("P-head", 18, True),
            ("дифф. VSH-декодер", 13, False),
        ],
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )
    add_arrow(s, 20.2, 10.15, width=0.7, height=0.4)

    # Box 5 — reconstructed P̂.
    add_rect(s, 21.1, 9.5, 4.0, 1.6, RGBColor(0xCC, 0xDD, 0xFF))
    add_text(
        s, 21.1, 9.5, 4.0, 1.6,
        text_blocks=[
            ("P̂ ∈ ℝ⁶⁴⁴⁴⁰", 16, True),
            ("(восстановленная мощность)", 13, False),
        ],
        anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
    )

    # Loss branches below boxes 3 and 5.
    add_text(
        s, 10.9, 11.4, 4.0, 0.6,
        text_blocks=[("↓", 18, True)],
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s, 10.9, 12.0, 4.0, 0.7,
        text_blocks=[("coef_mse: ‖ŷ − y_target‖²", 14, False)],
        align=PP_ALIGN.CENTER,
    )

    add_text(
        s, 21.1, 11.4, 4.0, 0.6,
        text_blocks=[("↓", 18, True)],
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s, 21.1, 12.0, 4.0, 0.7,
        text_blocks=[("power_loss: ‖P̂ − P_target‖²_w", 14, False)],
        align=PP_ALIGN.CENTER,
    )

    # Gradient-flow note (bottom).
    add_text(
        s, 1.2, 12.8, 24.5, 0.7,
        text_blocks=[("Градиенты текут справа налево через P-head: ∂L/∂η = (∂L/∂P̂) · (∂P̂/∂ŷ) · (∂ŷ/∂η).", 16, True)],
        color=RGBColor(0x55, 0x55, 0x55),
    )

    add_footer(s, 5, TOTAL, sw, sh)
    add_notes(
        s,
        "Бейзлайн — простой MLP пять на двести. На реальных "
        "данных лучший набор признаков — raw_plus_sh: сырая мощность плюс энергия по сферическим "
        "гармоникам. Сквозной пайплайн: вход (признаки P) → MLP → предсказанные packed-коэффициенты "
        "ŷ длины 1 020 → P-head, который через дифференцируемый VSH-декодер строит "
        "восстановленную мощность P̂ → две функции потерь (coef_mse на ŷ и power_loss на P̂). "
        "Градиенты текут через P-head — это позволяет обучать модель напрямую на физически "
        "осмысленную метрику."
    )


# ======================================================================
# SLIDE 6 — Метрики качества
# ======================================================================
def _slide_9_metrics_DRAFT(prs, blank, sw, sh) -> None:
    s = prs.slides.add_slide(blank)
    add_title(s, "Метрики качества")

    # Two columns — training losses (left), evaluation metrics (right).
    add_text(
        s, 0.8, 2.0, 12.5, 1.0,
        text_blocks=[("Функции потерь при обучении", 30, True)],
    )
    add_bullets(
        s, 0.8, 3.2, 12.5, 6.5,
        bullets=[
            (0, "coef_mse_loss — прямая MSE на упакованный вектор коэффициентов ŷ.", 22),
            (0, "power_loss — area-weighted MSE на P̂ через дифференцируемый VSH-декодер.", 22),
            (0, "Регуляризационные лоссы:", 22),
            (1, "physics_power_rank,", 20),
            (1, "rank_bin_p - назначаем ранг P̂ и считаем лосс над рангами,", 20),
            (1, "physics_power_mixed (coef_mse + λ · physics_power).", 20),
        ],
    )

    # Highlight: the p_rank finding.
    add_rect(s, 0.8, 9.9, 12.5, 3.5, RGBColor(0xFF, 0xF5, 0xCC))
    add_text(
        s, 1.0, 10.1, 12.1, 1.0,
        text_blocks=[("Ключевая находка:", 22, True)],
        color=RGBColor(0xC0, 0x10, 0x10),
    )
    add_bullets(
        s, 1.0, 11.0, 12.1, 2.4,
        bullets=[
            (0, "power_loss с rank_bin_p — основной регуляризатор: качественный скачок в скорости сходимости и итоговом качестве модели.", 20),
            (0, "Аблейшн: baseline vs baseline + p_rank — графики добавим к финальной версии.", 20),
        ],
    )

    # Right column — evaluation metrics + acceptance criteria.
    add_text(
        s, 14.0, 2.0, 12.5, 1.0,
        text_blocks=[("Метрики сравнения моделей", 30, True)],
    )
    add_bullets(
        s, 14.0, 3.2, 12.5, 4.5,
        bullets=[
            (0, "accuracy over bins — точность на P-точках, разбитых по бинам мощности.", 22),
            (0, "r² — per-P-point coefficient of determination.", 22),
            (0, "Обе метрики оцениваются над всеми точками поля P (per-pixel), затем берётся медиана/распределение.", 22),
        ],
    )

    add_text(
        s, 14.0, 8.0, 12.5, 1.0,
        text_blocks=[("Критерии приёмки модели:", 24, True)],
    )
    add_bullets(
        s, 14.0, 9.0, 12.5, 4.5,
        bullets=[
            (0, "median(r²) > 0.5 — выше тривиального предсказания среднего.", 22),
            (0, "accuracy > random — модель не угадывает.", 22),
            (0, "Оба критерия должны выполняться одновременно на val и holdout.", 22),
        ],
    )

    add_footer(s, 6, TOTAL, sw, sh)
    add_notes(
        s,
        "Слева — функции потерь, на которых обучаем модель. Главная "
        "находка работы: добавление power_loss с rank-based вариантом (p_rank) к coef_mse даёт "
        "качественный скачок в сходимости и итоговом качестве модели. Конкретный эксперимент — "
        "baseline против baseline + p_rank будет показан в финальной версии. Справа — метрики, "
        "по которым мы сравниваем модели: per-P-point accuracy и r². Критерии приёмки: "
        "медиана r² должна сходиться выше 0.5, accuracy — выше случайного предсказания, и "
        "оба условия — одновременно на валидационной и отложенной выборках. Финальные формулы "
        "и числа аблейшна я заменю до защиты."
    )



# ----------------------------------------------------------------------
# Result slides (8–11): metrics + plots for the four runs.
# Figures are taken from experiments/baseline/figures_real_augmented_*.
# Numerical metrics are read from experiments/baseline/S5_real_augmented_results_*.json
# (where available).
# ----------------------------------------------------------------------
import json as _json

_RUNS_DIR = ROOT / "experiments" / "baseline"


def _load_run_metrics(json_name: str) -> dict | None:
    p = _RUNS_DIR / json_name
    if not p.exists():
        return None
    try:
        return _json.load(open(p, encoding="utf-8"))
    except Exception:
        return None


def _fmt(v):
    if v is None:
        return "—"
    try:
        v = float(v)
    except Exception:
        return str(v)
    if abs(v) < 10:
        return f"{v:.3g}"
    return f"{v:.2e}"


def _metrics_rows(metrics: dict) -> list[tuple[str, str, str, str]]:
    """Return rows of (split, coef_amb, field_NRMSE_w, bin_acc) for the standard splits."""
    rows = []
    for split in ("train_aug", "val_real", "holdout_real"):
        rows.append((
            split,
            _fmt(metrics.get(f"report/{split}/coef_mse_amb_aware")),
            _fmt(metrics.get(f"report/{split}/field_nrmse_w")),
            _fmt(metrics.get(f"report/{split}/p_bin_accuracy")),
        ))
    return rows


def _result_slide(prs, blank, sw, sh, *, n: int, title: str, subtitle: str,
                  fig_dir: Path, r2_pdf: Path, bin_pdf: Path,
                  metrics_json: str | None) -> None:
    s = prs.slides.add_slide(blank)
    add_title(s, title, size=36)
    add_text(
        s, 0.8, 1.7, 25.5, 0.8,
        text_blocks=[(subtitle, 20, False)],
        color=RGBColor(0x55, 0x55, 0x55),
    )

    # Two plots side-by-side, full available height (метрики читаются прямо с графиков).
    img_top = 3.0
    img_h = 10.0

    add_text(s, 0.8, img_top - 0.7, 12.5, 0.6,
             text_blocks=[("r² distribution", 20, True)])
    if r2_pdf.exists():
        png = pdf_first_page_to_png(r2_pdf, zoom=2.0)
        add_image(s, png, 0.8, img_top, 12.5, img_h)
    else:
        add_rect(s, 0.8, img_top, 12.5, img_h, RGBColor(0xEE, 0xEE, 0xEE))
        add_text(s, 0.8, img_top, 12.5, img_h,
                 text_blocks=[("[r²: figure not found]", 16, False)],
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
                 color=RGBColor(0x88, 0x88, 0x88))

    add_text(s, 14.0, img_top - 0.7, 12.5, 0.6,
             text_blocks=[("accuracy over P-bins", 20, True)])
    if bin_pdf.exists():
        png = pdf_first_page_to_png(bin_pdf, zoom=2.0)
        add_image(s, png, 14.0, img_top, 12.5, img_h)
    else:
        add_rect(s, 14.0, img_top, 12.5, img_h, RGBColor(0xEE, 0xEE, 0xEE))
        add_text(s, 14.0, img_top, 12.5, img_h,
                 text_blocks=[("[bin-accuracy: figure not found]", 16, False)],
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
                 color=RGBColor(0x88, 0x88, 0x88))

    add_footer(s, n, TOTAL, sw, sh)


def _slide_run_best_small(prs, blank, sw, sh) -> None:
    fig = ROOT / "experiments" / "baseline" / "figures_real_augmented_best_small"
    _result_slide(
        prs, blank, sw, sh,
        n=7,
        title="Run «best_small»: mlp_3x200 + physics_power_rank + raw_plus_sh",
        subtitle="Малый прогон главного сетапа (raw_plus_sh + p_rank), 1 000 augmented samples, эпоха 10. JSON отчёт пока не дампился — только графики.",
        fig_dir=fig,
        r2_pdf=fig / "epoch_0010" / "r2_distribution.pdf",
        bin_pdf=fig / "epoch_0010" / "bin_accuracy_distribution.pdf",
        metrics_json=None,
    )
    add_notes(
        prs.slides[-1],
        "Результат №1 — best_small. Главный пайплайн (mlp_3x200 + physics_power_rank + "
        "raw_plus_sh), 1 000 аугментированных сэмплов, графики на эпохе 10. JSON-отчёт "
        "ещё не дампился, поэтому метрики на этом слайде нет — только распределения r² "
        "и accuracy по бинам мощности. Это контрольный мини-прогон перед полным сетапом."
    )


def _slide_run_coef(prs, blank, sw, sh) -> None:
    fig = ROOT / "experiments" / "baseline" / "figures_real_augmented_coef"
    _result_slide(
        prs, blank, sw, sh,
        n=8,
        title="Run «coef»: mlp_5x200 + coef_mse + raw_plus_sh",
        subtitle="Бейзлайн без physics_power: чистый coef_mse, 10 000 augmented samples, L=5.",
        fig_dir=fig,
        r2_pdf=fig / "r2_distribution.pdf",
        bin_pdf=fig / "bin_accuracy_distribution.pdf",
        metrics_json="S5_real_augmented_results_coef.json",
    )
    add_notes(
        prs.slides[-1],
        "Результат №2 — coef. Тот же feature-set (raw_plus_sh), но лосс — чистый coef_mse, "
        "без physics_power. Бейзлайн для аблейшна с p_rank. На реальных сплитах val_real "
        "и holdout_real coef_r² близок к нулю (≈ −0.01), а bin accuracy около 0.09 — "
        "практически на уровне случайного предсказания."
    )


def _slide_run_coef_raw_flat(prs, blank, sw, sh) -> None:
    fig = ROOT / "experiments" / "baseline" / "figures_real_augmented_coef_raw_flat"
    _result_slide(
        prs, blank, sw, sh,
        n=9,
        title="Run «coef_raw_flat»: mlp_5x200 + coef_mse + raw_flat",
        subtitle="Тот же coef_mse, но без spherical-harmonic канала: только raw P. Замер вклада sh-фичей.",
        fig_dir=fig,
        r2_pdf=fig / "r2_distribution.pdf",
        bin_pdf=fig / "bin_accuracy_distribution.pdf",
        metrics_json="S5_real_augmented_results_coef_raw_flat.json",
    )
    add_notes(
        prs.slides[-1],
        "Результат №3 — coef_raw_flat. Контроль на тип признаков: без sh-канала, только "
        "raw P. На coef_mse + 10 000 augmented разница с raw_plus_sh минимальна (val_real "
        "coef_amb ≈ 0.48 vs 0.48). Sh-канал на coef_mse не помогает — выигрыш виден только "
        "в комбинации с physics_power_rank."
    )


def _slide_run_physics(prs, blank, sw, sh) -> None:
    fig = ROOT / "experiments" / "baseline" / "figures_real_augmented_physics"
    _result_slide(
        prs, blank, sw, sh,
        n=10,
        title="Run «physics»: physics_power_rank без coef_mse",
        subtitle="Чистый physics_power_rank лосс (без coef_mse-компоненты), графики на эпохе 30. JSON отчёт пока не дампился.",
        fig_dir=fig,
        r2_pdf=fig / "epoch_0030" / "r2_distribution.pdf",
        bin_pdf=fig / "epoch_0030" / "bin_accuracy_distribution.pdf",
        metrics_json=None,
    )
    add_notes(
        prs.slides[-1],
        "Результат №4 — physics. Чистый physics_power_rank, без coef_mse. JSON-отчёт ещё "
        "не дампился, метрики не показываем; зато r² distribution и accuracy по бинам "
        "видны на 30-й эпохе. Это аблейшн в обратную сторону относительно coef-runs: "
        "проверяем, насколько rank-based лосс работает соло."
    )


def _slide_run_best(prs, blank, sw, sh) -> None:
    fig = ROOT / "experiments" / "baseline" / "figures_real_augmented_best"
    _result_slide(
        prs, blank, sw, sh,
        n=11,
        title="Run «best»: mlp_3x200 + physics_power_rank + raw_plus_sh (epoch 80)",
        subtitle="Главный full-эпоха-80 прогон того же сетапа, что best_small.",
        fig_dir=fig,
        r2_pdf=fig / "epoch_0080" / "r2_distribution.pdf",
        bin_pdf=fig / "epoch_0080" / "bin_accuracy_distribution.pdf",
        metrics_json=None,
    )
    add_notes(
        prs.slides[-1],
        "Результат №5 — best (full, эпоха 80). Тот же сетап, что best_small (mlp_3x200 + "
        "physics_power_rank + raw_plus_sh), доведённый до 80-й эпохи. r²-distribution и "
        "accuracy-over-bins — слева и справа. Карта поля на holdout_real — на следующем слайде."
    )


def _slide_holdout_field_grid(prs, blank, sw, sh) -> None:
    """Slide 12/13 — holdout_real field_comparison_grid across all 5 runs."""
    BASE = ROOT / "experiments" / "baseline"
    runs = [
        ("best (epoch 80)",
         BASE / "figures_real_augmented_best" / "epoch_0080" / "holdout_real" / "field_comparison_grid.pdf"),
        ("best_small (epoch 10)",
         BASE / "figures_real_augmented_best_small" / "epoch_0010" / "holdout_real" / "field_comparison_grid.pdf"),
        ("coef",
         BASE / "figures_real_augmented_coef" / "holdout_real" / "field_comparison_grid.pdf"),
        ("coef_raw_flat",
         BASE / "figures_real_augmented_coef_raw_flat" / "holdout_real" / "field_comparison_grid.pdf"),
        ("physics (epoch 30)",
         BASE / "figures_real_augmented_physics" / "epoch_0030" / "holdout_real" / "field_comparison_grid.pdf"),
    ]

    s = prs.slides.add_slide(blank)
    add_title(s, "Holdout_real: сетка P_pred vs P_true по всем 5 запускам", size=32)
    add_text(
        s, 0.8, 1.4, 25.5, 0.6,
        text_blocks=[("Слева для каждого ряда — название run-а; справа — pred / true / ratio для нескольких источников holdout.", 16, False)],
        color=RGBColor(0x55, 0x55, 0x55),
    )

    # Horizontal layout: 5 columns side-by-side, each with a header row above the plot.
    n_cols = len(runs)
    left0 = 0.4
    col_w = 5.16          # 5 * 5.16 + 4 * 0.05 ≈ 26.0
    col_gap = 0.05
    header_top = 2.0
    header_h = 0.7
    plot_top = header_top + header_h + 0.1
    plot_h = 14.0 - plot_top - 0.4  # leaves space above footer

    for i, (label, pdf) in enumerate(runs):
        x = left0 + i * (col_w + col_gap)
        # Header label.
        add_rect(s, x, header_top, col_w, header_h, RGBColor(0xF7, 0xF7, 0xF7))
        add_text(
            s, x, header_top, col_w, header_h,
            text_blocks=[(label, 14, True)],
            anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
        )
        # Plot beneath.
        if pdf.exists():
            png = pdf_first_page_to_png(pdf, zoom=2.0)
            add_image(s, png, x, plot_top, col_w, plot_h)
        else:
            add_rect(s, x, plot_top, col_w, plot_h, RGBColor(0xEE, 0xEE, 0xEE))
            add_text(
                s, x, plot_top, col_w, plot_h,
                text_blocks=[(f"[not found: {pdf.name}]", 12, False)],
                anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER,
                color=RGBColor(0x88, 0x88, 0x88),
            )

    add_footer(s, 12, TOTAL, sw, sh)
    add_notes(
        s,
        "Holdout_real — сетки сравнения предсказанной и истинной мощности для всех пяти "
        "запусков подряд: best (full, эпоха 80), best_small (эпоха 10), coef, coef_raw_flat "
        "и physics (эпоха 30). Каждая строка — один run; внутри каждой сетки слева P_pred, "
        "в середине P_true, справа карта отношения. Источники holdout — не видны модели при "
        "обучении. Сравнение позволяет видеть, насколько каждый сетап схватывает крупную "
        "угловую структуру на out-of-source-распределении."
    )


def _slide_conclusions(prs, blank, sw, sh) -> None:
    """Slide 13/13 — выводы и полезность результата."""
    s = prs.slides.add_slide(blank)
    add_title(s, "Выводы")

    # Left column — какие подходы применили и почему.
    add_text(
        s, 0.8, 2.0, 12.5, 1.0,
        text_blocks=[("Применённые подходы", 28, True)],
    )
    add_text(
        s, 0.8, 3.0, 12.5, 1.6,
        text_blocks=[(
            "Решение обратной задачи мультипольного анализа нетривиально: большой "
            "размер входа P, отсутствие единственного решения, дефицит обучающих "
            "данных. В работе использованы:",
            18, False
        )],
    )
    add_bullets(
        s, 0.8, 4.7, 12.5, 8.5,
        bullets=[
            (0, "Feature extraction — raw_flat, raw_plus_sh, FFT-radial, sh_power, power_pca, subsample_stride4.", 20),
            (0, "Полусинтетические данные — аугментации поверх реальных антенн (phi_roll, mode_dropout, additive_noise, coef_phase_rotation).", 20),
            (0, "Регуляризации разной природы — coef_mse, physics_power, p_rank, rank_bin_p, physics_power_mixed.", 20),
            (0, "Отдельная голова VSH-декодер — дифференцируемый форвард, через который текут градиенты от лоссов на P.", 20),
            (0, "Метрики двух типов — регрессионные (coef_mse, field_NRMSE_w, r²) и ranking (Spearman ρ, accuracy over P-bins).", 20),
        ],
    )

    # Right column — useful properties of the result.
    add_text(
        s, 14.0, 2.0, 12.5, 1.0,
        text_blocks=[("Полезность результата", 28, True)],
    )
    add_bullets(
        s, 14.0, 3.0, 12.5, 10.5,
        bullets=[
            (0, "VSH-декодер встроен в инференс: пользователь видит истинное и предсказанное P, может оценить расхождение и при необходимости откатиться на детерминированные методы.", 20),
            (0, "Бюджет обучения и инференса (на железе [Z]):", 20),
            (1, "обучение лучшей модели на одной антенне — [X] часов;", 18),
            (1, "инференс на одном сэмпле — [Y].", 18),
            (0, "Создан расширяемый фреймворк для дальнейших экспериментов в задаче мультипольного анализа.", 20),
        ],
    )

    # Closing line.
    add_text(
        s, 0.8, 13.6, 12.0, 1.0,
        text_blocks=[(f"Связь:  {SUBSTITUTIONS['FIO']}  ·  {SUBSTITUTIONS['EMAIL']}", 18, False)],
    )
    add_text(
        s, 14.5, 13.6, 12.0, 1.0,
        text_blocks=[("Спасибо. Готов к вопросам.", 22, True)],
        align=PP_ALIGN.RIGHT,
    )

    add_footer(s, 13, TOTAL, sw, sh)
    add_notes(
        s,
        "Финальный слайд. Слева — почему задача нетривиальна и какие пять подходов мы "
        "применили: feature extraction, полусинтетические данные через аугментации поверх "
        "реальных антенн, разные регуляризации, отдельная голова VSH-декодер для "
        "дифференцируемого форварда, и две группы метрик — регрессия и ранжирование. "
        "Справа — три пункта о полезности результата: пользователь может через VSH-декодер "
        "сверить P_pred с P_true и откатиться на детерминированные методы при недовольстве; "
        "обучение на одной антенне занимает [X] часов на [Z] и инференс одного сэмпла — [Y] "
        "(подставлю реальные значения); и создан фреймворк, на котором можно дальше ставить "
        "эксперименты в этой области."
    )


if __name__ == "__main__":
    out = build()
    print(f"Built: {out}")
    print(f"Size : {out.stat().st_size / 1024:.1f} KB")
